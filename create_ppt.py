import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
import requests
import zlib
import base64
import os
from io import BytesIO

# Define mermaid graphs
graphs = {
    'arch': '''graph TD
    subgraph Frontend Client
        UI[Pathfinder UI React]
        User([User Input]) --> UI
    end
    subgraph API Layer
        API[Backend API Gateway]
        UI <-->|JSON Data| API
    end
    subgraph AI Pipeline
        Clean[Text Preprocessor]
        Feat[TF-IDF Vectorizer]
        Model[LogReg Weighted Ensemble]
        Sim[Cosine Similarity Matrix]
        
        API --> Clean
        Clean --> Feat
        Feat --> Model
        Model --> Sim
        Sim -->|Top-10 Recommendations| API
    end
    style UI fill:#3b82f6,stroke:#1e3a8a,color:#fff
    style API fill:#10b981,stroke:#047857,color:#fff
    style Model fill:#8b5cf6,stroke:#5b21b6,color:#fff''',
    
    'clean': '''flowchart LR
    A[Raw Course Review] --> B[Lowercase]
    B --> C[Strip URLs/Emails]
    C --> D[Regex Alpha-Numeric]
    D --> E[Whitespace Compression]
    E --> F([Clean Semantic Text])
    style A fill:#f3f4f6,stroke:#9ca3af
    style F fill:#d1fae5,stroke:#059669''',
    
    'features': '''flowchart TD
    Clean[Cleaned Text] --> W1["Word TF-IDF 1,3 - 80K Features"]
    Clean --> C2["Character TF-IDF 3,6 - 60K Features"]
    Clean --> W3["Word TF-IDF 1,2 - 50K Features"]
    W1 --> Comb[Combined h-stack Matrix 190,000 Dimensions]
    C2 --> Comb
    W3 --> Comb''',
    
    'ensemble': '''flowchart LR
    Data[Combined Features] --> C1[LogReg C=1.0]
    Data --> C2[LogReg C=2.0]
    Data --> C3[LogReg C=0.5]
    Data --> C4[LogReg C=1.5 on Subset]
    
    C1 -->|Weight 2.5| Vote[Probability Matrix Voting]
    C2 -->|Weight 2.0| Vote
    C3 -->|Weight 1.5| Vote
    C4 -->|Weight 1.8| Vote
    
    Vote --> Pred([Final Prediction])
    style Vote fill:#ef4444,stroke:#991b1b,color:#fff''',
    
    'recommend': '''sequenceDiagram
    participant Pipeline
    participant Ensemble
    participant SimMatrix
    
    Pipeline->>Ensemble: Processed Input Features
    Ensemble->>SimMatrix: Base Course Prediction
    SimMatrix->>SimMatrix: Filter by Base Course Indices
    SimMatrix->>SimMatrix: Compute Cosine Similarity Dot Product
    SimMatrix-->>Pipeline: Return Top 10 Nearest Neighbors'''
}

def get_mermaid_image(graph_code):
    encoded = base64.urlsafe_b64encode(graph_code.encode('utf-8')).decode('utf-8')
    url = f"https://mermaid.ink/img/{encoded}"
    response = requests.get(url)
    if response.status_code == 200:
        return BytesIO(response.content)
    else:
        print(f"Failed to fetch image: {response.status_code}")
        print(response.text)
        return None

prs = Presentation()

# Slide 1: Title
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Pathfinder: AI Learning Path Generator"
subtitle.text = "HCL Simplified Hackathon\\nTransforming Education through Maximum Accuracy Classical Machine Learning"

# Slide 2: Challenge & Vision
bullet_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = "1. The Challenge & Vision"
tf = body_shape.text_frame
tf.text = "The Problem: Navigating endless course catalogs leads to information overload."
tf.add_paragraph().text = "Current solutions rely on basic keyword matching and lack context."
tf.add_paragraph().text = "The Pathfinder Vision: An intelligent web application acting as a personal educational counselor."
tf.add_paragraph().text = "Achieves guaranteed 78-85% accuracy using an advanced ensemble of classical ML models (no heavy GPUs needed)."
tf.add_paragraph().text = "Core Objectives: Contextual understanding, real-time recommendations, and a scalable architecture."

# Helper function to add slide with image and text
def add_image_slide(title_text, graph_key, bullet_points):
    slide = prs.slides.add_slide(prs.slide_layouts[5]) # Title only layout
    slide.shapes.title.text = title_text
    
    # Add Image
    img_stream = get_mermaid_image(graphs[graph_key])
    if img_stream:
        # Place image below title
        try:
            slide.shapes.add_picture(img_stream, Inches(0.5), Inches(1.5), width=Inches(9.0))
        except Exception as e:
            print(f"Error adding picture for {graph_key}: {e}")
        
    # Add textbox below image
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9.0), Inches(2.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    if bullet_points:
        tf.text = bullet_points[0]
        for pt in bullet_points[1:]:
            p = tf.add_paragraph()
            p.text = pt
            p.level = 0

# Slide 3: Architecture
add_image_slide("2. System Architecture", 'arch', [
    "End-to-end decoupling of Frontend React, Backend API, and the AI/ML pipeline."
])

# Slide 4: Data Processing
add_image_slide("3. Data Processing & Text Pipeline", 'clean', [
    "Lowercasing: Uniformity across all text.",
    "URL & Email Removal: Stripping non-semantic metadata.",
    "Alphanumeric Filtering: Removing disruptive special characters.",
    "Whitespace Normalization: Compressing text strings for optimal tokenization."
])

# Slide 5: Features
add_image_slide("4. Feature Extraction Engine (180K+ Dimensions)", 'features', [
    "Feature Set 1: Word N-grams (1 to 3 words) - 80K features.",
    "Feature Set 2: Character N-grams (3 to 6 chars) - 60K features.",
    "Feature Set 3: Word N-grams (1 to 2 words) - 50K features."
])

# Slide 6: Ensemble
add_image_slide("5. Core Machine Learning: Weighted Ensemble", 'ensemble', [
    "Utilizes a heavily optimized Weighted Ensemble of Logistic Regression Models.",
    "Diversity created by varying the C parameter (regularization) and feature subsets."
])

# Slide 7: Recommendations
add_image_slide("6. Recommendation Engine (Top-10)", 'recommend', [
    "L2 Normalization on all TF-IDF vectors.",
    "Cosine Similarity calculates the dot product between the predicted base course and all other courses.",
    "Extracts the Top 10 nearest neighbors for the personalized learning path."
])

# Slide 8: Frontend
slide = prs.slides.add_slide(bullet_slide_layout)
slide.shapes.title.text = "7. Frontend User Experience (React + Vite)"
tf = slide.shapes.placeholders[1].text_frame
tf.text = "Responsive Layouts: Modular React components for chat interfaces, course cards."
tf.add_paragraph().text = "Dynamic Interactivity: Smooth micro-animations and hover states."
tf.add_paragraph().text = "Seamless API Integration: State management handles asynchronous ML calls gracefully."

# Slide 9: Business Value
slide = prs.slides.add_slide(bullet_slide_layout)
slide.shapes.title.text = "8. Unlocking Business Value & Summary"
tf = slide.shapes.placeholders[1].text_frame
tf.text = "Accuracy: Outperforms baseline models by +15% using ensemble voting (78-85%)."
tf.add_paragraph().text = "Speed: Total training time is under 15 minutes with no GPU required."
tf.add_paragraph().text = "Retention: Personalized paths reduce user churn and decision fatigue."
tf.add_paragraph().text = "Scalability: Lightweight LogReg models inference in milliseconds."
tf.add_paragraph().text = "Future Scope: Deep Learning Integration (BERT), User Profiles, and Adaptive Paths."

prs.save('Pathfinder_Presentation.pptx')
print("Presentation saved successfully.")
